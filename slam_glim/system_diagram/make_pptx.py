#!/usr/bin/env python3
"""Generate GLIM Map Building Automation System Architecture presentation."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

DIR         = os.path.dirname(os.path.abspath(__file__))
TEMPLATE    = os.path.join(DIR, 'aaa.pptx')
OUTPUT      = os.path.join(DIR, 'GLIM_Map_Building_Automation_System_Architecture.pptx')
IMG_MINIMAL = os.path.join(DIR, 'GLIM_Map_Building_Automation_System_MINIMAL.png')
IMG_FULL    = os.path.join(DIR, 'GLIM_Map_Building_Automation_System.png')

# ── Load template (inherits master, theme, fonts, colors) ────────────────────
prs = Presentation(TEMPLATE)
N_ORIG = len(prs.slides)

COVER   = prs.slide_layouts[11]   # Cover
SECTION = prs.slide_layouts[13]   # Section Title_1
CONTENT = prs.slide_layouts[14]   # Content_A-1  (title / subtitle / content area)

# ── Helpers ───────────────────────────────────────────────────────────────────
def phs(slide):
    return {ph.placeholder_format.idx: ph for ph in slide.placeholders}

def set_text(ph, text):
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].text = text

def add_bullets(ph, items):
    """items: list of str (level 0) or (str, level) tuples."""
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in items:
        text  = item if isinstance(item, str) else item[0]
        level = 0    if isinstance(item, str) else item[1]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text  = text
        p.level = level
        first   = False

def cover(title, name, role, date):
    s = prs.slides.add_slide(COVER)
    p = phs(s)
    set_text(p[0],  title)
    set_text(p[13], name)
    set_text(p[14], role)
    set_text(p[15], date)

def section(title):
    s = prs.slides.add_slide(SECTION)
    set_text(phs(s)[0], title)

def content(title, subtitle, bullets=None, image=None, img_full_width=True):
    s  = prs.slides.add_slide(CONTENT)
    p  = phs(s)
    set_text(p[0],  title)
    set_text(p[11], subtitle)

    # Content area coords (from layout inspection)
    CA_L = Inches(0.57);  CA_T = Inches(2.39)
    CA_W = Inches(12.19); CA_H = Inches(4.14)

    if bullets and 12 in p:
        add_bullets(p[12], bullets)

    if image:
        if img_full_width:
            s.shapes.add_picture(image, CA_L, CA_T, width=CA_W)
        else:
            pic = s.shapes.add_picture(image, CA_L, CA_T, height=CA_H)
            # centre horizontally
            pic.left = int((Inches(13.33) - pic.width) / 2)
    return s

# ── Slides ────────────────────────────────────────────────────────────────────

# 1 · Cover
cover(
    "GLIM Map Building\nAutomation System",
    "Abdullah AlMusalami",
    "Senior Robotics Engineer",
    "May 2026",
)

# 2 · Overview section break
section("System Overview")

# 3 · Minimal architecture diagram
content(
    "GLIM Map Building Automation System",
    "3-Phase System Architecture",
    image=IMG_MINIMAL,
    img_full_width=True,
)

# 4 · Phase 1 section break
section("Phase 1  ·  Docker Environment")

# 5 · Phase 1 details
content(
    "Phase 1  ·  Docker Environment",
    "Environment Setup",
    bullets=[
        "Base image: nvidia/cuda:13.2 + Ubuntu 22.04",
        "ROS2 Jazzy Desktop",
        "GLIM — GPU-accelerated 3D SLAM (koide3 PPA)",
        "All dependencies pre-installed: GTSAM 4.3, Iridescence, GLFW, METIS",
        "GTSAM version conflict resolved inside the container — zero manual steps",
        "ROS2 bag files mounted at runtime  ( -v /path/to/bags:/bags )",
        "GPU passthrough:  --gpus all  (requires nvidia-container-toolkit)",
    ],
)

# 6 · Phase 2 section break
section("Phase 2  ·  Interactive CLI Wizard")

# 7 · Full pipeline diagram
content(
    "Phase 2  ·  Interactive CLI Wizard",
    "Full Pipeline:    ros2 run slam_glim main",
    image=IMG_FULL,
    img_full_width=False,
)

# 8 · Steps 1 & 1.5
content(
    "Phase 2  ·  CLI Wizard",
    "Step 1: Bag Preparation    +    Step 1.5: Sensor Configuration",
    bullets=[
        "Step 1  —  Bag Preparation:",
        ("User provides the bag file path", 1),
        ("ROS1 .bag?  →  auto-convert to ROS2 .db3 (rosbags-convert)", 1),
        "Step 1.5  —  Sensor Configuration:",
        ("User edits ONE YAML file:   config/user_config.yaml", 1),
        ("Only 3 fields required:", 1),
        ("imu_topic:    /imu/gravity", 2),
        ("lidar_topic:  /velodyne_points", 2),
        ("T_lidar_imu:  [tx, ty, tz, qx, qy, qz, qw]", 2),
        ("System auto-fills all 6 GLIM JSON config files — user never edits JSON", 1),
    ],
)

# 9 · Step 2
content(
    "Phase 2  ·  CLI Wizard",
    "Step 2: Sensor Validation",
    bullets=[
        "Runs GLIM's built-in validator_node against the bag",
        ("Checks: IMU rate, LiDAR rate, topic existence, gravity vector detection", 1),
        "Prints a human-readable report to the terminal",
        "Human reviews the report  →  fix config/user_config.yaml if needed",
        "Catches misconfigured topics before wasting a full SLAM run",
        "Press Enter when satisfied to proceed to SLAM",
    ],
)

# 10 · Steps 3-4
content(
    "Phase 2  ·  CLI Wizard",
    "Steps 3–4: GLIM SLAM Execution",
    bullets=[
        "Pipeline selection  (asked every run):",
        ("[1] GPU         —  VGICP_GPU + CUDA   (default, RMC robot)", 1),
        ("[2] CPU         —  VGICP CPU, no GPU required", 1),
        ("[3] LiDAR-only  —  CT-ICP, no IMU required", 1),
        "Config review: shows IMU topic, LiDAR topic, T_lidar_imu — confirm before SLAM starts",
        "Mode selection:",
        ("[1] Live Viewer   —  glim_rosnode + ros2 bag play + Iridescence 3D viewer", 1),
        ("[2] Fast Offline  —  glim_rosbag, no viewer, faster than real-time", 1),
        "Dump auto-saved  →  maps/<bag_name>_<YYYYMMDD_HHMMSS>/",
    ],
)

# 11 · Steps 5-6
content(
    "Phase 2  ·  CLI Wizard",
    "Steps 5–6: Map Refinement & Merging",
    bullets=[
        "Step 5  —  Edit / Modify Map:",
        ("Opens GLIM map_editor GUI with the current dump", 1),
        ("Add or remove loop closures, re-optimize the pose graph", 1),
        ("Loop: edit more?  →  repeat until satisfied", 1),
        "Step 6  —  Merge Maps:",
        ("User provides a second dump directory", 1),
        ("map_editor loads both dumps  →  user aligns and merges in the GUI", 1),
        ("Merged dump path saved for the export step", 1),
    ],
)

# 12 · Steps 7-8
content(
    "Phase 2  ·  CLI Wizard",
    "Steps 7–8: Map Export & 2D Map Generation",
    bullets=[
        "Step 7  —  Export Map  (fully automated, no GUI):",
        ("Reads dump binary: points_compact.bin (float32 XYZ) + data.txt (T_world_origin)", 1),
        ("Transforms all submaps to world frame  →  voxel filter 0.05 m", 1),
        ("Writes ASCII PCD v0.7  →  maps/<run>/map.pcd", 1),
        ("No PLY step — dump binary is parsed directly in Python", 1),
        "Step 8  —  2D Occupancy Map:",
        ("Ground detection: z-histogram from 1st percentile, find floor peak", 1),
        ("Height filter: floor+0.3 m  to  floor+3.0 m", 1),
        ("Output: map_2d.pgm + map_2d.yaml  (nav2_map_server compatible)", 1),
    ],
)

# 13 · Step 9
content(
    "Phase 2  ·  CLI Wizard",
    "Step 9: POI Registration",
    bullets=[
        "RViz2 opens automatically with the 3D PCD map displayed",
        "Select the Publish Point tool in the RViz2 toolbar (crosshair icon)",
        "Click any point on the 3D cloud",
        "Red sphere + text label appear at the clicked location",
        "POI saved automatically to YAML in  maps/<run>/POI_Poses/",
        "YAML format:",
        ("pois: [ { id, name, timestamp, frame_id, position, orientation } ]", 1),
        ("orientation: identity quaternion — compatible with ROS2 nav goals", 1),
    ],
)

# 14 · Phase 3 section break
section("Phase 3  ·  Web Dashboard")

# 15 · Phase 3 features
content(
    "Phase 3  ·  Web Dashboard",
    "Browser-Based Interface  (Coming Soon)",
    bullets=[
        "Replaces CLI wizard with a web page — same pipeline, zero terminal interaction",
        "3D map viewer in the browser  (Three.js / Potree)",
        "POI editor: click-to-place on the 3D map in the browser",
        "Live robot tracking: WebSocket stream of robot pose overlaid on the global map",
        "Map library: browse all saved map runs with metadata and thumbnails",
        "OEM export: download PCD + 2D map + POI YAML for distribution",
        "Architecture:",
        ("Flask / FastAPI backend  →  calls the same step modules as the CLI", 1),
        ("Design constraint: all step modules accept decisions as parameters (no input())", 1),
    ],
)

# 16 · Implementation plan
content(
    "Implementation Plan",
    "Phase 2 — Implementation Order",
    bullets=[
        "1.   export_map.py           —  Binary dump → PCD  (core unique logic, first)",
        "2.   convert_map.py          —  PCD → 2D occupancy grid",
        "3.   pcd_publisher.py  +  poi_selector.py   (copy + adapt from glim_test)",
        "4.   view_map.launch.py  +  view_map.rviz",
        "5.   config_setup.py         —  user_config.yaml → GLIM JSON files",
        "6.   slam_runner.py          —  GLIM launch: pipeline select + two modes",
        "7.   validator.py            —  Sensor validation subprocess + report",
        "8.   bag_convert.py          —  ROS1 bag conversion wrapper",
        "9.   map_editor_launcher.py  —  Map editing / merging subprocess",
        "10.  main.py                 —  CLI wizard orchestrator (wires all steps)",
    ],
)

# ── Remove original template slides (reverse order keeps indices valid) ───────
REL_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

def delete_slide(prs, idx):
    sld_id_lst = prs.slides._sldIdLst
    sId  = sld_id_lst[idx]
    rId  = sId.get(f'{{{REL_NS}}}id')
    sld_id_lst.remove(sId)
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass

for i in range(N_ORIG - 1, -1, -1):
    delete_slide(prs, i)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved:  {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
